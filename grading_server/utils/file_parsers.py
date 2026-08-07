"""
parsers.py
File parsers for PDF, DOCX, and plain text submissions.
PDFs/DOCX now surface tables and describe embedded figures via Gemini vision
so grading can credit charts/images that have no extractable text.
"""

from __future__ import annotations

import io
import re
import zipfile


def parse_submission(file_bytes: bytes, content_type: str, api_key: str | None = None) -> str:
    """
    Parse a submission file into plain text.
    Supports: application/pdf, application/vnd.openxmlformats-officedocument.wordprocessingml.document, text/*
    """
    ct = content_type.lower()

    if "pdf" in ct:
        return _parse_pdf(file_bytes, api_key=api_key)
    elif "wordprocessingml" in ct or "docx" in ct:
        return _parse_docx(file_bytes, api_key=api_key)
    elif "presentationml" in ct or "pptx" in ct:
        return _parse_pptx(file_bytes)
    elif "ms-powerpoint" in ct or "ppt" in ct:
        return "This is a legacy .ppt file; content could not be extracted for AI grading. Please review manually."
    elif "text" in ct:
        return file_bytes.decode("utf-8", errors="replace")
    else:
        try:
            return file_bytes.decode("utf-8", errors="replace")
        except Exception:
            raise ValueError(f"Unsupported content type: {content_type}")


def _describe_visual(image_bytes: bytes, mime_type: str = "image/png", api_key: str | None = None) -> str:
    """Describe a chart/image/table screenshot for inclusion in submission text."""
    if not image_bytes or len(image_bytes) < 200:
        return ""
    try:
        from google import genai
        from google.genai import types
        from grading_server.utils.gemini_client import _genai_client, _with_quota_retry

        client = genai.Client(api_key=api_key) if api_key else _genai_client
        prompt = (
            "You are helping grade a student submission. Describe this page/figure for a text-only grader.\n"
            "Include: chart/table/image type, title/caption if visible, axes/legend labels, "
            "key values or trends, and whether it is a real data visualization (not decorative).\n"
            "Be concrete and concise (max ~180 words). Do not say the figure is missing."
        )

        def _call():
            return client.models.generate_content(
                model="models/gemini-2.5-flash",
                contents=[
                    prompt,
                    types.Part.from_bytes(data=image_bytes, mime_type=mime_type),
                ],
            )

        response = _with_quota_retry(_call, max_attempts=3)
        return (response.text or "").strip()
    except Exception as e:
        print(f"[VISUAL EXTRACT] Error: {e}")
        return ""


def _page_needs_vision(page, text: str) -> bool:
    """Heuristic: page likely contains charts/images that text extraction misses."""
    try:
        images = page.get_images(full=True) or []
    except Exception:
        images = []
    if images:
        return True

    # Vector charts (Excel/PowerPoint paste) often appear as drawings, not images
    try:
        drawings = page.get_drawings() or []
    except Exception:
        drawings = []
    if len(drawings) >= 12:
        return True

    lower = (text or "").lower()
    visual_cues = (
        "data visual", "visualization", "figure", "fig.", "chart", "graph",
        "plot", "histogram", "scatter", "bar chart", "table", "appendix",
    )
    if any(cue in lower for cue in visual_cues) and len(drawings) >= 3:
        return True
    return False


def _extract_pdf_tables(file_bytes: bytes) -> dict[int, list[str]]:
    """page_index -> markdown-ish table strings via pdfplumber."""
    out: dict[int, list[str]] = {}
    try:
        import pdfplumber
    except Exception:
        return out

    try:
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                tables = []
                try:
                    raw_tables = page.extract_tables() or []
                except Exception:
                    raw_tables = []
                for t_idx, table in enumerate(raw_tables, start=1):
                    if not table:
                        continue
                    rows = []
                    for row in table:
                        cells = [(" ".join(str(c).split()) if c is not None else "") for c in row]
                        if any(cells):
                            rows.append(" | ".join(cells))
                    if len(rows) >= 2:
                        tables.append(f"[TABLE {t_idx}]\n" + "\n".join(rows))
                if tables:
                    out[i] = tables
    except Exception as e:
        print(f"[PDF TABLES] Error: {e}")
    return out


def _parse_pdf(file_bytes: bytes, api_key: str | None = None) -> str:
    """
    Extract PDF text + tables, and describe pages that contain figures/charts
    (vision) so graders can see visualizations that have no OCR text.
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    table_map = _extract_pdf_tables(file_bytes)
    pages: list[str] = []
    vision_budget = 8  # cap Gemini page renders per submission

    for i, page in enumerate(doc):
        page_no = i + 1
        chunks: list[str] = [f"=== Page {page_no} ==="]
        text = (page.get_text("text") or "").strip()
        if text:
            chunks.append(text)

        for table_block in table_map.get(i, []):
            chunks.append(table_block)

        if vision_budget > 0 and _page_needs_vision(page, text):
            try:
                # Render page so charts/images are visible to Gemini
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                png_bytes = pix.tobytes("png")
                description = _describe_visual(png_bytes, "image/png", api_key=api_key)
                if description:
                    chunks.append(
                        f"[FIGURE / VISUAL on page {page_no} - treat as present in the submission]\n"
                        f"{description}"
                    )
                    vision_budget -= 1
                else:
                    chunks.append(
                        f"[FIGURE / VISUAL DETECTED on page {page_no}: "
                        f"an embedded chart or image is present but could not be fully described. "
                        f"Do NOT claim visuals are missing on this page.]"
                    )
                    vision_budget -= 1
            except Exception as e:
                print(f"[PDF VISION] page {page_no}: {e}")
                chunks.append(
                    f"[FIGURE / VISUAL DETECTED on page {page_no}: embedded visual present.]"
                )

        if len(chunks) > 1:
            pages.append("\n\n".join(chunks))

    doc.close()
    return "\n\n".join(pages)


def _docx_image_bytes(file_bytes: bytes, max_images: int = 6) -> list[bytes]:
    """Pull embedded image binaries from a DOCX zip."""
    images: list[bytes] = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            names = [
                n for n in zf.namelist()
                if n.startswith("word/media/")
                and re.search(r"\.(png|jpe?g|gif|webp|bmp)$", n, re.I)
            ]
            names.sort()
            for name in names[:max_images]:
                data = zf.read(name)
                if data and len(data) >= 500:
                    images.append(data)
    except Exception as e:
        print(f"[DOCX IMAGES] Error: {e}")
    return images


def _parse_docx(file_bytes: bytes, api_key: str | None = None) -> str:
    """Extract text from a DOCX (paragraphs, tables, headers/footers, figures)."""
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(io.BytesIO(file_bytes))
    parts: list[str] = []

    def _para_text(p: Paragraph) -> str:
        return (p.text or "").strip()

    def _table_text(table: Table) -> str:
        rows: list[str] = []
        for row in table.rows:
            cells = [" ".join(c.text.split()) for c in row.cells if c.text and c.text.strip()]
            seen: set[str] = set()
            unique: list[str] = []
            for cell in cells:
                if cell not in seen:
                    seen.add(cell)
                    unique.append(cell)
            if unique:
                rows.append(" | ".join(unique))
        return "\n".join(rows)

    # Body: walk document order so table content isn't lost
    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
        if tag == "p":
            text = _para_text(Paragraph(block, doc))
            if text:
                parts.append(text)
        elif tag == "tbl":
            text = _table_text(Table(block, doc))
            if text:
                parts.append(f"[TABLE]\n{text}")

    # Headers / footers often hold titles or student names
    for section in doc.sections:
        header_candidates = []
        footer_candidates = []
        try:
            header_candidates = [section.header, section.first_page_header, section.even_page_header]
            footer_candidates = [section.footer, section.first_page_footer, section.even_page_footer]
        except Exception:
            try:
                header_candidates = [section.header]
                footer_candidates = [section.footer]
            except Exception:
                pass
        for header in header_candidates:
            if header is None:
                continue
            try:
                for p in header.paragraphs:
                    text = _para_text(p)
                    if text:
                        parts.append(text)
            except Exception:
                pass
        for footer in footer_candidates:
            if footer is None:
                continue
            try:
                for p in footer.paragraphs:
                    text = _para_text(p)
                    if text:
                        parts.append(text)
            except Exception:
                pass

    # Fallback if XML walk missed content
    if not parts:
        parts = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            text = _table_text(table)
            if text:
                parts.append(f"[TABLE]\n{text}")

    # Describe embedded images/charts (common for pasted Excel charts)
    for idx, img in enumerate(_docx_image_bytes(file_bytes), start=1):
        mime = "image/png"
        if img[:3] == b"\xff\xd8\xff":
            mime = "image/jpeg"
        description = _describe_visual(img, mime, api_key=api_key)
        if description:
            parts.append(
                f"[FIGURE / VISUAL {idx} - treat as present in the submission]\n{description}"
            )
        else:
            parts.append(
                f"[FIGURE / VISUAL {idx} DETECTED: an embedded image/chart is present in the document. "
                f"Do NOT claim visuals are missing.]"
            )

    return "\n\n".join(parts)


def _parse_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if getattr(shape, "has_table", False):
                try:
                    table = shape.table
                    for row in table.rows:
                        cells = [" ".join(c.text.split()) for c in row.cells if c.text.strip()]
                        if cells:
                            lines.append(" | ".join(cells))
                except Exception:
                    pass
        # Note pictures so graders don't claim "no images"
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            pic_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            if pic_count:
                lines.append(
                    f"[FIGURE / VISUAL: {pic_count} image(s)/chart(s) embedded on this slide — "
                    f"treat as present in the submission]"
                )
        except Exception:
            pass
        lines.append("")
    return "\n".join(lines)
