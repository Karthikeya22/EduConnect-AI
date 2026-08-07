import sys
import os
import tempfile
import asyncio
import time
import traceback
from typing import TypedDict, Literal
from pydantic import BaseModel

class QualityCriterion(BaseModel):
    criterion_id: str
    title: str
    description: str
    max_score: int = 10
    dimension: Literal["content", "structure"] = "structure"

class DynamicCriteria(BaseModel):
    criteria: list[QualityCriterion]
app_dir = os.path.dirname(os.path.abspath(__file__))
workspace_dir = os.path.dirname(app_dir)
workspace_tmp = os.path.join(workspace_dir, "tmp")
os.makedirs(workspace_tmp, exist_ok=True)
os.environ["TEMP"] = workspace_tmp
os.environ["TMP"] = workspace_tmp
tempfile.tempdir = workspace_tmp

from flask import Flask, request, jsonify
from flask_cors import CORS
from langgraph.graph import StateGraph, END

# Use verified local utils
from grading_server.utils.supabase_client import get_supabase
from grading_server.utils.gemini_client import (
    call_gemini, call_gemini_text, embed_texts
)
from grading_server.utils.file_parsers import parse_submission, _parse_pdf, _parse_docx, _parse_pptx
from grading_server.config import GEMINI_LITE_MODEL, GEMINI_GRADING_MODEL, FLASK_PORT, CANVAS_BASE_URL
from grading_server.models import CriterionVerdict, CritiqueResult, GradingOutput

# Windows event loop policy fix
if sys.platform == "win32":
    try:
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    except Exception:
        pass

app = Flask(__name__)
CORS(app)

class GradingState(TypedDict):
    assignment_id: str
    course_id: str | None
    student_id: str
    submission_text: str
    chunks: list[str]
    chunk_embeddings: list[list[float]]
    submission_concepts: list[str]
    rubric_criteria: list[dict]
    relevant_chunks: list[dict]
    exemplars: dict  # criterion_id -> list of exemplar dicts
    verdicts: list[dict]
    critiques: list[dict]
    final_output: dict
    loop_count: int
    force_flag_for_human: bool
    custom_gemini_key: str | None

# Node 1: Ingest (Sync)
async def run_ingest(state: GradingState) -> dict:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    txt = state["submission_text"]
    chunks = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50).split_text(txt) or [txt[:1000]]
    api_key = state.get("custom_gemini_key")
    # Sync call in thread
    c_res = await asyncio.to_thread(call_gemini, f"Extract concepts as a JSON object with a 'concepts' key (array of strings): {txt[:2000]}", model=GEMINI_LITE_MODEL, api_key=api_key)
    e_res = await asyncio.to_thread(embed_texts, chunks, api_key=api_key)
    
    concepts = []
    if isinstance(c_res, list):
        concepts = c_res
    elif isinstance(c_res, dict):
        concepts = c_res.get("concepts", [])
        
    return {"chunks": chunks, "chunk_embeddings": e_res, "submission_concepts": concepts}

# Node 2: Retrieve — Semantic Vector Search + Exemplars
async def run_retrieve(state: GradingState) -> dict:
    import numpy as np

    def _db():
        sb = get_supabase()
        aid = state["assignment_id"]
        embeddings = state.get("chunk_embeddings", [])

        try:
            # 1. Fetch rubric criteria for this assignment
            rubric = sb.table("rubric_criteria").select("*").eq("assignment_id", aid).execute().data or []

            # 3. Fetch exemplars for few-shot grading context
            exemplars_raw = sb.table("exemplars").select("*").eq("assignment_id", aid).execute().data or []
            # Group exemplars by criterion_id
            exemplars_by_criterion = {}
            for ex in exemplars_raw:
                cid = ex.get("criterion_id", "")
                if cid not in exemplars_by_criterion:
                    exemplars_by_criterion[cid] = []
                exemplars_by_criterion[cid].append(ex)

            # 4. Generate Adaptive Quality Evaluation Criteria
            if rubric:
                try:
                    from grading_server.utils.gemini_client import call_gemini
                    
                    quality_prompt = (
                        "Based on the following assignment rubric, generate exactly 2 core Quality & Academic Integrity evaluation criteria tailored to this specific assignment type. "
                        "IMPORTANT: The Academic Integrity criterion MUST be highly lenient and give the student the benefit of the doubt by default. It should ONLY flag genuine, obvious integrity concerns (e.g., direct blatant plagiarism) and must NOT penalize stylistic or structural similarities. "
                        "For example, if it's a research paper, focus on Methodology Soundness and Citation Integrity. If it's programming, focus on Code Quality and Best Practices. "
                        "These criteria will be used to automatically grade the submission.\n\n"
                        f"Existing Rubric Titles: {[r.get('title') for r in rubric]}\n"
                    )
                    api_key = state.get("custom_gemini_key")
                    quality_res = call_gemini(quality_prompt, model=GEMINI_LITE_MODEL, response_schema=DynamicCriteria, api_key=api_key)
                    new_criteria = quality_res.get("criteria", [])
                    for nc in new_criteria:
                        nc_dict = dict(nc)
                        nc_dict["assignment_id"] = aid
                        nc_dict["title"] = f"{nc_dict.get('title', 'Criterion')} (AI-Generated)"
                        nc_dict["ai_generated"] = True
                        rubric.append(nc_dict)
                    print(f"[RAG] Injected {len(new_criteria)} adaptive quality criteria.")
                except Exception as e:
                    print(f"[RAG] Failed to inject dynamic quality criteria: {e}")

            print(f"[RAG] Retrieved: {len(rubric)} criteria, {len(exemplars_raw)} exemplars")
            return rubric, exemplars_by_criterion
        except Exception as e:
            print(f"[RAG] Retrieve error: {e}")
            return [], [], {}

    r, ex = await asyncio.to_thread(_db)
    return {"rubric_criteria": r, "exemplars": ex, "relevant_chunks": []}

from pydantic import BaseModel
from typing import Literal

class LLMGradingPayload(BaseModel):
    score: int
    justification: str
    evidence_anchor: str
    status: Literal['full', 'partial', 'missing']
    missing_keywords: list[str] = []
    supporting_materials: list[str] = []
    covered_concepts: list[str] = []
    required_concepts: list[str] = []

class DecomposedCriterion(BaseModel):
    sub_queries: list[str]

async def _decompose_criterion(criterion_title: str, criterion_desc: str, api_key: str | None = None) -> list[str]:
    prompt = f"""
You are an expert grading assistant. Your task is to decompose the following grading rubric criterion into 2 to 4 specific, verifiable sub-questions. These sub-questions will be used as search queries to retrieve relevant course materials.
Focus on the core factual, conceptual, or structural requirements of the criterion.

Rubric Criterion Title: {criterion_title}
Rubric Criterion Description: {criterion_desc}

Decompose this into 2-4 sub-queries. Return them as a JSON list of strings.
"""
    try:
        from grading_server.utils.gemini_client import call_gemini
        res = await asyncio.to_thread(call_gemini, prompt, model=GEMINI_LITE_MODEL, response_schema=DecomposedCriterion, api_key=api_key)
        return res.get("sub_queries", [])
    except Exception as e:
        print(f"[Decompose Error] {e}")
        return [f"{criterion_title}: {criterion_desc}"]

async def _multihop_retrieve(
    sub_queries: list[str],
    assignment_id: str,
    api_key: str | None = None,
    course_id: str | None = None,
) -> list[dict]:
    if not sub_queries:
        return []
    from grading_server.utils.gemini_client import embed_texts
    try:
        embeddings = await asyncio.to_thread(embed_texts, sub_queries, api_key=api_key)
    except Exception as e:
        print(f"[Embed Error] {e}")
        return []
    
    def _db():
        sb = get_supabase()
        all_chunks = []
        for vec in embeddings:
            if not vec: continue
            try:
                params = {
                    "query_embedding": vec,
                    "match_count": 5,
                    "filter_assignment_id": assignment_id,
                }
                if course_id:
                    params["filter_course_id"] = course_id
                try:
                    chunks = sb.rpc("match_course_material_chunks", params).execute().data or []
                except Exception as rpc_err:
                    # Live DB may still have the 3-arg RPC (migration not applied).
                    # Retry without filter_course_id so retrieval still works.
                    err_text = str(rpc_err)
                    if course_id and ("PGRST202" in err_text or "filter_course_id" in err_text or "Could not find the function" in err_text):
                        legacy = {
                            "query_embedding": vec,
                            "match_count": 5,
                            "filter_assignment_id": assignment_id,
                        }
                        chunks = sb.rpc("match_course_material_chunks", legacy).execute().data or []
                    else:
                        raise
                all_chunks.extend(chunks)
            except Exception as e:
                print(f"[Vector Search Error] {e}")
        if not all_chunks:
            # Last-resort: return any stored chunks for this assignment (no similarity)
            try:
                fallback = (
                    sb.table("course_material_chunks")
                    .select("id, chunk_text, source_title, chunk_index")
                    .eq("assignment_id", assignment_id)
                    .limit(8)
                    .execute()
                    .data
                    or []
                )
                for row in fallback:
                    row.setdefault("similarity", 0.0)
                all_chunks.extend(fallback)
            except Exception as e:
                print(f"[Vector Search Fallback Error] {e}")
        return all_chunks
    
    chunks = await asyncio.to_thread(_db)
    seen = set()
    unique_chunks = []
    for c in chunks:
        cid = c.get("id")
        if cid not in seen:
            seen.add(cid)
            unique_chunks.append(c)
    return unique_chunks

# Parallel Grade Helper — Context-Aware with Semantic RAG + Exemplars
async def _gp(c, ctx, txt, exemplar_context, previous_critique=None, api_key=None):
    criterion_title = c.get('title', c.get('criterion_id', 'Unknown'))
    criterion_desc = c.get('description', '')
    max_score = c['max_score']

    # Build exemplar section if available
    exemplar_section = ""
    if exemplar_context:
        exemplar_section = "\n--- EXEMPLAR EXAMPLES (for calibration) ---\n"
        for ex in exemplar_context[:3]:  # limit to top 3 exemplars
            exemplar_section += (
                f"  Score: {ex.get('score', '?')}/{ex.get('max_score', '?')} — "
                f"Excerpt: {str(ex.get('submission_text', ''))[:300]}\n"
            )
        exemplar_section += "--- END EXEMPLARS ---\n\n"

    if previous_critique:
        exemplar_section += f"=== PREVIOUS CRITIQUE (Self-Correction Needed) ===\n"
        exemplar_section += f"The previous grading attempt was flagged. Critique: {previous_critique.get('reasoning', 'Unknown')}\n"
        exemplar_section += f"Please revise your justification and score to address this critique.\n\n"

    # Heuristic for dimension if defaulted to 'content'
    dimension = c.get('dimension', 'content')
    title_lower = criterion_title.lower()
    if any(kw in title_lower for kw in ['structure', 'format', 'grammar', 'organization', 'presentation', 'style', 'mechanic', 'writing']):
        dimension = 'structure'

    if dimension == 'structure':
        scoring_philosophy = (
            f"=== SCORING PHILOSOPHY (STRUCTURE) ===\n"
            f"- Be fair, constructive, and evidence-based.\n"
            f"- PRIMARY EVIDENCE: Grade ONLY from the STUDENT SUBMISSION. evidence_anchor MUST be a verbatim quote from the student text.\n"
            f"- If the student demonstrated the required structure, extract a direct student quote as 'evidence_anchor'.\n"
            f"- VISUALS/TABLES: [FIGURE / VISUAL ...] and [TABLE ...] blocks are real student content — credit them; never say images/tables are missing when present.\n"
            f"- If no valid student quote exists, set 'evidence_anchor' to 'not found'. Do NOT quote course slides/lectures as evidence.\n"
            f"- Course materials are OPTIONAL REFERENCE only (definitions, examples). Do NOT require the student to match lecture examples or slide filenames.\n"
            f"- Focus on logical organization (clear sections, flow, transitions) and cohesion.\n"
            f"- EXPLICITLY IGNORE visual formatting details you cannot see.\n"
            f"- Do not penalize minor grammar issues or non-native phrasing as long as the structure and meaning are clear.\n"
            f"- 'missing_keywords' should list rubric requirements the student failed to include (not unrelated lecture topics).\n"
            f"- 'supporting_materials' should list course Source titles only if they truly informed this judgment; otherwise use [].\n"
            f"- CONCEPT GRAPH COVERAGE: Derive 'required_concepts' from the RUBRIC first. 'covered_concepts' must come from the student text.\n\n"
        )
    else:
        scoring_philosophy = (
            f"=== SCORING PHILOSOPHY ===\n"
            f"- Be fair, constructive, and evidence-based.\n"
            f"- PRIMARY EVIDENCE: Grade against the RUBRIC using the STUDENT SUBMISSION. evidence_anchor MUST be a verbatim quote from the student text.\n"
            f"- If the student addressed the criterion (including synonyms/paraphrase), give appropriate credit.\n"
            f"- VISUALS/TABLES: Blocks marked [FIGURE / VISUAL ...] or [TABLE ...] are real content from the student's file (charts, images, tables). Treat them as present evidence. Do NOT claim 'no visualizations/images/tables' when those markers exist.\n"
            f"- If no valid student quote exists, set 'evidence_anchor' to 'not found'. Never invent anchors or cite slide decks as student evidence.\n"
            f"- Course materials are OPTIONAL REFERENCE for terminology/context — NOT mandatory ground truth. Do NOT zero a criterion because lecture slides omit that section, or because the student's domain/example differs from a slide example.\n"
            f"- Prefer the rubric description over unrelated retrieved lecture chunks.\n"
            f"- 'missing_keywords' should list rubric concepts the student missed — not topics that only appear in course slides.\n"
            f"- 'supporting_materials' should list Source titles only when actually used; otherwise [].\n"
            f"- CONCEPT GRAPH COVERAGE: Derive 'required_concepts' from the RUBRIC first. 'covered_concepts' must come from the student text.\n\n"
        )

    prompt = (
        f"You are an expert academic grader. Grade the student's submission against ONE specific rubric criterion.\n\n"
        f"=== RUBRIC CRITERION ===\n"
        f"Title: {criterion_title}\n"
        f"Description: {criterion_desc}\n"
        f"Max Score: {max_score}\n\n"
        f"{scoring_philosophy}"
        f"=== OPTIONAL COURSE CONTEXT (reference only — not required) ===\n{ctx}\n\n"
        f"{exemplar_section}"
        f"=== STUDENT SUBMISSION (grade this) ===\n{txt[:150000]}\n"
    )
    try:
        r = await asyncio.to_thread(call_gemini, prompt, model=GEMINI_GRADING_MODEL, response_schema=LLMGradingPayload, api_key=api_key)
    except Exception as e:
        print(f"[GRADE] Error grading criterion {c.get('criterion_id')}: {e}")
        r = {"score": 0, "status": "missing", "justification": f"Grading error: {str(e)}", "evidence_anchor": "not found", "missing_keywords": []}

    # Clamp score to max
    raw_score = int(r.get('score', 0))
    clamped_score = max(0, min(raw_score, max_score))
    evidence_anchor = str(r.get('evidence_anchor', '') or '').strip()

    # Reject anchors that are not actually present in the student submission
    if evidence_anchor and evidence_anchor.lower() not in ("not found", "n/a", "none"):
        norm_sub = " ".join(txt.lower().split())
        norm_anchor = " ".join(evidence_anchor.lower().split())
        # Allow short fuzzy match: full anchor or first ~12 words
        anchor_probe = norm_anchor if len(norm_anchor) <= 180 else " ".join(norm_anchor.split()[:12])
        if anchor_probe and anchor_probe not in norm_sub:
            evidence_anchor = "not found"

    v = {
        "criterion_id": c['criterion_id'],
        "criterion_name": c.get('title', c['criterion_id']),
        "dimension": dimension,
        "score": clamped_score,
        "max_score": max_score,
        "status": str(r.get('status', 'partial')),
        "justification": str(r.get('justification', '')),
        "missing_keywords": r.get('missing_keywords', []) or [],
        "evidence_anchor": evidence_anchor or "not found",
        "supporting_materials": r.get('supporting_materials', []) or [],
        "covered_concepts": r.get('covered_concepts', []) or [],
        "required_concepts": r.get('required_concepts', []) or [],
        "ai_generated": c.get("ai_generated", False)
    }
    return v

# Node 3: Grade (Concurrent) — Semantic RAG + Exemplar-Aware
async def run_grade(state: GradingState) -> dict:
    rubric = state.get("rubric_criteria", [])
    exemplars = state.get("exemplars", {})
    txt = state["submission_text"]
    aid = state["assignment_id"]

    if not rubric:
        return {"verdicts": [], "relevant_chunks": []}

    # Map previous critiques by criterion_id if this is a retry loop
    prev_critiques = {c["criterion_id"]: c for c in state.get("critiques", []) if c.get("flag_for_human")}

    async def _process_criterion(c):
        cid = c.get("criterion_id", "")
        criterion_title = c.get('title', cid)
        criterion_desc = c.get('description', '')
        
        api_key = state.get("custom_gemini_key")
        # 1. Decompose
        sub_queries = await _decompose_criterion(criterion_title, criterion_desc, api_key=api_key)
        # 2. Multi-Hop Retrieve
        retrieved_chunks = await _multihop_retrieve(
            sub_queries,
            aid,
            api_key=api_key,
            course_id=state.get("course_id"),
        )
        
        # Build context from semantically retrieved course material
        ctx = "\n---\n".join([
            f"[Source: {ch.get('source_title', 'Unknown Source').split(' (chunk')[0].strip()}] (Chunk {ch.get('chunk_index', '?')}) (similarity: {round(ch.get('similarity', 0), 2)})\n{ch.get('chunk_text', '')}"
            for ch in retrieved_chunks[:10]  # Take top 10 distinct chunks for this criterion
        ])
        if not ctx:
            ctx = "(No course material retrieved. Grade using the rubric and student submission only.)"
        
        criterion_exemplars = exemplars.get(cid, []) if isinstance(exemplars, dict) else []
        pc = prev_critiques.get(cid)
        
        # 3. Grade
        verdict = await _gp(c, ctx, txt, criterion_exemplars, previous_critique=pc, api_key=api_key)
        return verdict, retrieved_chunks

    tasks = [_process_criterion(c) for c in rubric]
    results = list(await asyncio.gather(*tasks))
    
    verdicts = [res[0] for res in results]
    all_chunks = []
    seen_chunks = set()
    for res in results:
        for ch in res[1]:
            cid = ch.get("id")
            if cid not in seen_chunks:
                seen_chunks.add(cid)
                all_chunks.append(ch)

    # Do NOT downgrade full credit for weak course-slide similarity — open-ended
    # assignments often diverge from lecture examples. Flag only when the model
    # claimed full credit with no student evidence anchor.
    force_flag_for_human = False
    for v in verdicts:
        anchor = (v.get("evidence_anchor") or "").strip().lower()
        if v.get("status") == "full" and anchor in ("", "not found", "n/a", "none"):
            v["status"] = "partial"
            v["ungrounded_full"] = True
            v["justification"] = (
                (v.get("justification") or "")
                + "\n\nDowngraded: full credit claimed without a verifiable quote from the student submission."
            ).strip()
            force_flag_for_human = True

    return {"verdicts": verdicts, "loop_count": state.get("loop_count", 0) + 1, "force_flag_for_human": force_flag_for_human, "relevant_chunks": all_chunks}

# Critique model for structured LLM output
class LLMCritiquePayload(BaseModel):
    is_supported: bool
    is_justified: bool
    confidence: float
    reasoning: str

# Parallel Critique Helper — Actually calls the LLM
async def _cp(v, api_key=None):
    criterion_title = v.get("criterion_name", v.get("criterion_id", "Unknown"))
    prompt = (
        f"You are a quality assurance reviewer for AI-generated grades. Review this grading verdict:\n\n"
        f"Criterion: {criterion_title}\n"
        f"Score Given: {v.get('score', 0)}/{v.get('max_score', 0)}\n"
        f"Status: {v.get('status', 'unknown')}\n"
        f"Justification: {v.get('justification', '')[:500]}\n"
        f"Evidence Anchor: {v.get('evidence_anchor', '')[:300]}\n\n"
        f"Evaluate:\n"
        f"1. is_supported: Is the justification grounded in the cited evidence anchor? (true/false)\n"
        f"2. is_justified: Is the score consistent with the justification? (e.g. a 'missing' status shouldn't have a high score)\n"
        f"3. confidence: How confident are you in this verdict? (0.0 to 1.0)\n"
        f"4. reasoning: Brief explanation of your assessment.\n"
    )
    try:
        r = await asyncio.to_thread(call_gemini, prompt, model=GEMINI_LITE_MODEL, response_schema=LLMCritiquePayload, api_key=api_key)
        confidence = float(r.get('confidence', 0.5))
        return CritiqueResult(
            criterion_id=v["criterion_id"],
            is_supported=bool(r.get('is_supported', True)),
            is_justified=bool(r.get('is_justified', True)),
            confidence=max(0.0, min(1.0, confidence)),
            flag_for_human=confidence < 0.65,
        ).model_dump()
    except Exception as e:
        print(f"[CRITIQUE] Error critiquing {v.get('criterion_id')}: {e}")
        return CritiqueResult(
            criterion_id=v["criterion_id"],
            is_supported=True, is_justified=True,
            confidence=0.5, flag_for_human=True,
        ).model_dump()

# Node 4: Critique (Concurrent) — Real LLM Critique
async def run_critique(state: GradingState) -> dict:
    v = state.get("verdicts", [])
    api_key = state.get("custom_gemini_key")
    if not v: return {"critiques": []}
    return {"critiques": list(await asyncio.gather(*[_cp(item, api_key=api_key) for item in v]))}

# Node 5: Output — with topic radar, misconception hint, and concepts for Knowledge Graph
async def run_output(state: GradingState) -> dict:
    verts = [CriterionVerdict(**v) for v in state.get("verdicts", [])]
    critiques = [CritiqueResult(**c) for c in state.get("critiques", [])]

    # Separate content vs structure scores, excluding AI-generated criteria
    content_verts = [v for v in verts if v.dimension == "content" and not v.ai_generated]
    structure_verts = [v for v in verts if v.dimension == "structure" and not v.ai_generated]

    content_score = sum(v.score for v in content_verts)
    content_max = sum(v.max_score for v in content_verts)
    structure_score = sum(v.score for v in structure_verts)
    structure_max = sum(v.max_score for v in structure_verts)
    total = content_score + structure_score
    total_max = content_max + structure_max

    # Build topic mastery radar (percentage per criterion)
    topic_mastery_radar = {}
    for v in verts:
        if v.max_score > 0:
            if getattr(v, "required_concepts", None):
                depth_score = len(v.covered_concepts) / len(v.required_concepts)
            else:
                depth_score = v.score / v.max_score
            topic_mastery_radar[v.criterion_name or v.criterion_id] = round(depth_score, 2)

    # Aggregate missing keywords for misconception hint
    all_missing = []
    for v in verts:
        all_missing.extend(v.missing_keywords)
    misconception_hint = ""
    if all_missing:
        unique_missing = list(dict.fromkeys(all_missing))[:10]  # deduplicate, limit to 10
        misconception_hint = f"The submission may be missing key concepts: {', '.join(unique_missing)}."

    # Check if rubric was available
    rubric_missing = len(state.get("rubric_criteria", [])) == 0
    if rubric_missing:
        misconception_hint = "No rubric criteria found for this assignment. Please run Ingest Materials."

    # Extract referenced materials
    relevant_chunks = state.get("relevant_chunks", [])
    referenced_materials = []
    for ch in relevant_chunks:
        src = ch.get("source_title")
        if src:
            # Strip chunk suffix if present e.g. " (chunk 1/11)"
            base_src = src.split(" (chunk")[0].strip()
            if base_src not in referenced_materials:
                referenced_materials.append(base_src)
    referenced_materials = referenced_materials[:5]

    # Calculate RAG coverage from retrieved course-material chunks.
    # Prior threshold (0.65 + needing 5 hits) was unrealistically strict for
    # embedding cosine scores, so nearly every paper showed LOW.
    sims = sorted(
        (float(ch.get("similarity") or 0) for ch in relevant_chunks),
        reverse=True,
    )
    strong_match_count = sum(1 for s in sims if s >= 0.45)
    usable_count = sum(1 for s in sims if s >= 0.28)
    top_avg = (sum(sims[:5]) / min(5, len(sims))) if sims else 0.0

    if not sims or usable_count == 0:
        rag_coverage_level = "NONE"
    elif strong_match_count >= 3 or top_avg >= 0.50:
        rag_coverage_level = "HIGH"
    elif strong_match_count >= 1 or usable_count >= 3 or top_avg >= 0.35:
        rag_coverage_level = "MEDIUM"
    else:
        rag_coverage_level = "LOW"

    # Human review: require real uncertainty — not a single noisy critique flag.
    overall_confidence = sum(c.confidence for c in critiques) / len(critiques) if critiques else 0.9
    flagged_count = sum(1 for c in critiques if c.flag_for_human)
    majority_flagged = flagged_count >= max(1, (len(critiques) + 1) // 2) if critiques else False
    any_flagged = (
        state.get("force_flag_for_human", False)
        or overall_confidence < 0.65
        or majority_flagged
    )

    out = GradingOutput(
        assignment_id=state["assignment_id"],
        student_id=state["student_id"],
        content_score=content_score,
        content_max=content_max,
        structure_score=structure_score,
        structure_max=structure_max,
        total=total,
        total_max=total_max,
        criteria_verdicts=[v.model_dump() for v in verts],
        topic_mastery_radar=topic_mastery_radar,
        misconception_hint=misconception_hint,
        flag_for_human=any_flagged,
        overall_confidence=round(overall_confidence, 2),
        critiques=state.get("critiques", []),
        referenced_materials=referenced_materials,
        rag_coverage_level=rag_coverage_level,
        strong_match_count=strong_match_count,
    )
    res = out.model_dump()
    res["rubric_missing"] = rubric_missing
    # Include submission concepts for Knowledge Graph updates on the frontend
    res["submission_concepts"] = state.get("submission_concepts", [])

    # Save to grading_results table
    try:
        def _save():
            sb = get_supabase()
            aid = state["assignment_id"]
            sid = state["student_id"]
            existing = sb.table("grading_results").select("id").eq("assignment_id", aid).eq("student_id", sid).execute().data
            if existing:
                row_id = existing[0]["id"]
                sb.table("grading_results").update({
                    "result_json": res
                }).eq("id", row_id).execute()
            else:
                sb.table("grading_results").insert({
                    "assignment_id": aid,
                    "student_id": sid,
                    "result_json": res
                }).execute()
        await asyncio.to_thread(_save)
    except Exception as e:
        print(f"[OUTPUT] Failed to save grading result: {e}")

    return {"final_output": res}

def should_continue(state: GradingState):
    loop_count = state.get("loop_count", 0)
    critiques = state.get("critiques", [])
    if loop_count < 2 and any(c.get("flag_for_human", False) for c in critiques):
        print(f"[RAG] Critique flagged issues. Self-correcting (Loop {loop_count})...")
        return "grade"
    return "output"

def get_graph():
    g = StateGraph(GradingState)
    g.add_node("ingest", run_ingest)
    g.add_node("retrieve", run_retrieve)
    g.add_node("grade", run_grade)
    g.add_node("critique", run_critique)
    g.add_node("output", run_output)
    g.set_entry_point("ingest")
    g.add_edge("ingest", "retrieve")
    g.add_edge("retrieve", "grade")
    g.add_edge("grade", "critique")
    g.add_conditional_edges("critique", should_continue, {"grade": "grade", "output": "output"})
    g.add_edge("output", END)
    return g.compile()

@app.route("/api/grade", methods=["POST"])
async def grade_submission():
    s = time.time()
    try:
        if request.is_json:
            data = request.get_json(force=True)
        else:
            data = request.form.to_dict()
            
        # Extract submission text (priority: text field > uploaded file)
        submission_text = data.get("submission_text", "").strip()
        custom_gemini_key = request.headers.get("X-Gemini-Api-Key")
        
        extracted_parts = []
        if request.files:
            for uploaded_file in request.files.values():
                filename = (uploaded_file.filename or "").lower()
                file_bytes = uploaded_file.read()
                
                if file_bytes:
                    part = ""
                    if filename.endswith(".pdf"):
                        part = _extract_pdf(file_bytes, api_key=custom_gemini_key)
                    elif filename.endswith(".docx"):
                        part = _extract_docx(file_bytes, api_key=custom_gemini_key)
                    elif filename.endswith(".pptx"):
                        part = _extract_pptx(file_bytes)
                    elif filename.endswith(".ppt"):
                        part = "This is a legacy .ppt file; content could not be extracted for AI grading. Please review manually."
                    elif filename.endswith(".txt"):
                        part = file_bytes.decode("utf-8", errors="ignore")
                    elif filename.endswith(".ipynb"):
                        part = _extract_ipynb(file_bytes)
                    elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                        mime_type = "image/png"
                        if filename.endswith((".jpg", ".jpeg")): mime_type = "image/jpeg"
                        elif filename.endswith(".webp"): mime_type = "image/webp"
                        part = _extract_image(file_bytes, mime_type, api_key=custom_gemini_key)
                    if part and part.strip():
                        extracted_parts.append(part.strip())
        
        # Download from file_urls if provided (authenticated Canvas download)
        file_urls = data.get("file_urls", [])
        canvas_token = (
            data.get("canvas_token")
            or request.headers.get("X-Canvas-Token")
            or request.form.get("canvas_token")
        )

        def _is_placeholder_body(text: str) -> bool:
            t = (text or "").strip()
            if not t:
                return True
            lower = t.lower()
            return (
                lower.startswith("submission type:")
                or lower == "no submission content."
                or lower.startswith("no submission content")
            )

        # Ignore Canvas placeholder bodies when real files are attached
        if file_urls and _is_placeholder_body(submission_text):
            submission_text = ""

        download_errors: list[str] = []
        if file_urls:
            for file_info in file_urls:
                filename = ""
                try:
                    url = file_info.get("url") if isinstance(file_info, dict) else None
                    filename = (file_info.get("filename") or "").lower() if isinstance(file_info, dict) else ""
                    if not url:
                        continue
                    # Fall back to URL path for extension when filename is odd/missing
                    if not any(filename.endswith(ext) for ext in (
                        ".pdf", ".docx", ".pptx", ".ppt", ".txt", ".ipynb", ".png", ".jpg", ".jpeg", ".webp"
                    )):
                        from urllib.parse import urlparse, unquote
                        path_name = unquote(urlparse(url).path).lower()
                        filename = path_name.split("/")[-1] or filename

                    file_bytes = await asyncio.to_thread(_download_canvas_bytes, url, canvas_token)

                    if file_bytes:
                        # OOXML magic: zip header — prefer content sniff when extension ambiguous
                        is_ooxml = file_bytes[:2] == b"PK"

                        def _extract_one() -> str:
                            if filename.endswith(".pdf") or file_bytes[:4] == b"%PDF":
                                return _extract_pdf(file_bytes, api_key=custom_gemini_key)
                            if filename.endswith(".docx") or (is_ooxml and "word" in filename):
                                return _extract_docx(file_bytes, api_key=custom_gemini_key)
                            if filename.endswith(".pptx") or (is_ooxml and "ppt" in filename):
                                return _extract_pptx(file_bytes)
                            if filename.endswith(".ppt"):
                                return "This is a legacy .ppt file; content could not be extracted for AI grading. Please review manually."
                            if filename.endswith(".txt"):
                                return file_bytes.decode("utf-8", errors="ignore")
                            if filename.endswith(".ipynb"):
                                return _extract_ipynb(file_bytes)
                            if filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
                                mime_type = "image/png"
                                if filename.endswith((".jpg", ".jpeg")): mime_type = "image/jpeg"
                                elif filename.endswith(".webp"): mime_type = "image/webp"
                                return _extract_image(file_bytes, mime_type, api_key=custom_gemini_key)
                            if is_ooxml:
                                try:
                                    return _extract_docx(file_bytes, api_key=custom_gemini_key)
                                except Exception:
                                    return _extract_pptx(file_bytes)
                            return ""

                        extracted = await asyncio.to_thread(_extract_one)
                        if not extracted and not (
                            filename.endswith((".pdf", ".docx", ".pptx", ".ppt", ".txt", ".ipynb", ".png", ".jpg", ".jpeg", ".webp"))
                            or is_ooxml
                            or file_bytes[:4] == b"%PDF"
                        ):
                            download_errors.append(f"Unsupported file type: {filename or url}")
                            continue

                        if extracted and len(extracted.strip()) >= 40:
                            extracted_parts.append(extracted.strip())
                        elif extracted and extracted.strip():
                            # Keep short extracts but warn — better than discarding
                            extracted_parts.append(extracted.strip())
                            print(f"[WARN] Short extract ({len(extracted.strip())} chars) from {filename}")
                        else:
                            download_errors.append(f"{filename or 'attachment'}: extracted empty text")
                except Exception as e:
                    print(f"Failed to download or extract {filename} from URL: {e}")
                    download_errors.append(f"{filename or 'attachment'}: {e}")

        if extracted_parts:
            submission_text = (submission_text + "\n\n" + "\n\n".join(extracted_parts)).strip()

        # Placeholder bodies like "Submission type: online_upload" are not real content
        placeholder_only = _is_placeholder_body(submission_text)
        if file_urls and not extracted_parts:
            detail = "; ".join(download_errors) if download_errors else "Could not extract text from attached files"
            return jsonify({
                "error": f"Could not read submission file(s). {detail}",
                "details": detail,
            }), 400
        if file_urls and extracted_parts and len(submission_text) < 40:
            return jsonify({
                "error": "Submission file was downloaded but almost no text was extracted. The document may be image-only or corrupted.",
                "details": f"chars={len(submission_text)}; errors={download_errors}",
            }), 400
        
        if not submission_text or placeholder_only:
            return jsonify({"error": "No submission text or valid file provided"}), 400

        print(f"[DEBUG] Extracted submission text length for {data.get('student_id')}: {len(submission_text)} characters")
        if "[FIGURE" in submission_text or "[TABLE" in submission_text:
            print(f"[DEBUG] Visual/table markers present for {data.get('student_id')}")

        final = await get_graph().ainvoke({
            "assignment_id": data.get("assignment_id"), 
            "course_id": data.get("course_id") or request.form.get("course_id"),
            "student_id": data.get("student_id"), 
            "submission_text": submission_text,
            "custom_gemini_key": custom_gemini_key
        })
        
        res = final["final_output"]
        res["_pipeline_duration_seconds"] = round(time.time() - s, 2)
        res["extracted_submission_text"] = submission_text
        return jsonify(res), 200
    except Exception as e:
        traceback.print_exc()
        err_msg = str(e)
        if "expected 768 dimensions" in err_msg or "expected 3072 dimensions" in err_msg:
             return jsonify({"error": "DATABASE_DIMENSION_MISMATCH", "details": "Please run the setup_ingestion_tables.sql in Supabase SQL Editor to update your vector dimensions to 3072."}), 500
        if "PERMISSION_DENIED" in err_msg or "API key not valid" in err_msg or "API key was reported as leaked" in err_msg:
             return jsonify({"error": "Invalid Gemini API key. Update GEMINI_API_KEY or your custom key in settings.", "details": err_msg}), 401
        if "RESOURCE_EXHAUSTED" in err_msg or "exceeded your current quota" in err_msg or "429" in err_msg:
             return jsonify({"error": "Gemini API quota exceeded. Wait ~30 seconds and retry one student.", "details": err_msg}), 429
        if "cannot schedule new futures after shutdown" in err_msg:
             return jsonify({"error": "Grading server restarted mid-request. Please click Grade again.", "details": err_msg}), 503
        # Surface a short actionable message (not just a vague internal error)
        short = err_msg[:240] if err_msg else "Unknown error"
        return jsonify({"error": f"Grading failed: {short}", "details": err_msg}), 500

@app.route("/api/health")
def health(): return jsonify({"status":"ok"}), 200


# ══════════════════════════════════════════════════════════════════════════════
# Phase 2 — Ingestion Endpoints (Professor-triggered)
# ══════════════════════════════════════════════════════════════════════════════

from grading_server.ingest_service import ingest_course_material, delete_assignment_data

@app.route("/api/ingest", methods=["POST"])
async def ingest_json():
    """
    POST /api/ingest
    Accepts JSON body with course_material_text, rubric_criteria, and exemplars.
    Chunks, embeds, and upserts everything into Supabase.
    """
    try:
        data = request.get_json(force=True)
        assignment_id = data.get("assignment_id")
        if not assignment_id:
            return jsonify({"error": "assignment_id is required"}), 400

        course_material_text = data.get("course_material_text", "")
        rubric_criteria = data.get("rubric_criteria", [])
        exemplars = data.get("exemplars", [])

        if not course_material_text and not rubric_criteria:
            return jsonify({"error": "At least course_material_text or rubric_criteria must be provided"}), 400

        custom_gemini_key = request.headers.get("X-Gemini-Api-Key")

        result = await ingest_course_material(
            assignment_id=assignment_id,
            course_material_text=course_material_text,
            rubric_criteria=rubric_criteria,
            exemplars=exemplars if exemplars else None,
            api_key=custom_gemini_key,
        )
        return jsonify(result), 200

    except Exception:
        traceback.print_exc()
        return jsonify({"error": "Ingestion failed"}), 500


@app.route("/api/ingest/file", methods=["POST"])
async def ingest_file():
    """
    POST /api/ingest/file
    Accepts multipart form upload with: file (PDF/DOCX), assignment_id,
    rubric_criteria (JSON string), exemplars (JSON string, optional).
    Extracts text from the file, then runs the same ingestion pipeline.
    """
    import json as _json
    try:
        assignment_id = request.form.get("assignment_id")
        if not assignment_id:
            return jsonify({"error": "assignment_id is required"}), 400

        # Parse rubric_criteria from JSON string
        rubric_criteria_raw = request.form.get("rubric_criteria", "[]")
        try:
            rubric_criteria = _json.loads(rubric_criteria_raw)
        except _json.JSONDecodeError:
            return jsonify({"error": "rubric_criteria must be a valid JSON string"}), 400

        # Parse exemplars from JSON string (optional)
        exemplars_raw = request.form.get("exemplars", "[]")
        try:
            exemplars = _json.loads(exemplars_raw)
        except _json.JSONDecodeError:
            exemplars = []

        # Extract file
        if "file" not in request.files:
            return jsonify({"error": "No file uploaded. Send a PDF or DOCX file in the 'file' field."}), 400

        uploaded_file = request.files["file"]
        filename = (uploaded_file.filename or "").lower()
        file_bytes = uploaded_file.read()

        if not file_bytes:
            return jsonify({"error": "Uploaded file is empty"}), 400

        # Determine file type and extract text
        if filename.endswith(".pdf"):
            course_material_text = _extract_pdf(file_bytes)
        elif filename.endswith(".docx"):
            course_material_text = _extract_docx(file_bytes)
        elif filename.endswith(".pptx"):
            course_material_text = _extract_pptx(file_bytes)
        elif filename.endswith(".ppt"):
            course_material_text = "This is a legacy .ppt file; content could not be extracted for AI grading. Please review manually."
        elif filename.endswith(".txt"):
            course_material_text = file_bytes.decode("utf-8", errors="ignore")
        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
            mime_type = "image/png"
            if filename.endswith((".jpg", ".jpeg")): mime_type = "image/jpeg"
            elif filename.endswith(".webp"): mime_type = "image/webp"
            course_material_text = _extract_image(file_bytes, mime_type)
        else:
            return jsonify({
                "error": f"Unsupported file type: '{filename}'."
            }), 400

        if not course_material_text.strip():
            return jsonify({"error": "Could not extract any text from the uploaded file"}), 400

        custom_gemini_key = request.headers.get("X-Gemini-Api-Key")

        result = await ingest_course_material(
            assignment_id=assignment_id,
            course_material_text=course_material_text,
            rubric_criteria=rubric_criteria,
            exemplars=exemplars if exemplars else None,
            source_name=filename,
            api_key=custom_gemini_key,
        )
        return jsonify(result), 200

    except Exception:
        traceback.print_exc()
        return jsonify({"error": "File ingestion failed"}), 500

import ipaddress
import requests
from urllib.parse import urlparse

# Only Canvas-hosted files may be fetched server-side; the download is
# authenticated with the caller's Canvas token, so an arbitrary URL would let a
# client turn this route into an SSRF proxy.
MAX_CANVAS_DOWNLOAD_BYTES = 25 * 1024 * 1024


def _canvas_url_rejection(file_url: str) -> str | None:
    """Return a rejection reason for non-Canvas URLs, or None when allowed."""
    parsed = urlparse(file_url)
    if parsed.scheme not in ("http", "https"):
        return "Only http(s) URLs can be ingested"

    host = (parsed.hostname or "").lower().rstrip(".")
    if not host:
        return "URL has no host"

    if host in ("localhost", "localhost.localdomain", "metadata.google.internal"):
        return "Refusing to fetch a local address"

    try:
        ip = ipaddress.ip_address(host)
    except ValueError:
        ip = None
    if ip is not None and (
        ip.is_private or ip.is_loopback or ip.is_link_local or ip.is_reserved or ip.is_multicast
    ):
        return "Refusing to fetch a private or link-local address"

    allowed_hosts = {"instructure.com", "inscloudgate.net", "instructuremedia.com"}
    configured = urlparse(CANVAS_BASE_URL).hostname
    if configured:
        allowed_hosts.add(configured.lower())

    if host.startswith("canvas."):
        return None
    for allowed in allowed_hosts:
        if host == allowed or host.endswith(f".{allowed}"):
            return None
    return f"Host '{host}' is not a recognized Canvas host"


def _download_canvas_bytes(file_url: str, canvas_token: str | None = None) -> bytes:
    """Download a Canvas-hosted file with auth, size limit, and SSRF checks."""
    rejection = _canvas_url_rejection(str(file_url))
    if rejection:
        raise ValueError(rejection)

    headers = {"User-Agent": "EduConnect-AI/1.0"}
    if canvas_token:
        headers["Authorization"] = f"Bearer {canvas_token}"

    with requests.get(file_url, headers=headers, timeout=30, stream=True, allow_redirects=True) as r:
        if r.status_code != 200:
            raise ValueError(f"Failed to download from Canvas: {r.status_code}")

        declared_length = r.headers.get("Content-Length")
        if declared_length and declared_length.isdigit() and int(declared_length) > MAX_CANVAS_DOWNLOAD_BYTES:
            raise ValueError("File is too large (limit 25MB)")

        parts: list[bytes] = []
        total = 0
        for part in r.iter_content(chunk_size=65536):
            if not part:
                continue
            total += len(part)
            if total > MAX_CANVAS_DOWNLOAD_BYTES:
                raise ValueError("File is too large (limit 25MB)")
            parts.append(part)
        return b"".join(parts)


@app.route("/api/proxy-file", methods=["GET"])
def proxy_file():
    """
    Browser-safe proxy for Canvas file URLs (avoids CORS Failed to fetch in Local Reader).
    GET /api/proxy-file?url=...
    Header: X-Canvas-Token (optional)
    """
    file_url = request.args.get("url")
    if not file_url:
        return jsonify({"error": "url is required"}), 400

    canvas_token = request.headers.get("X-Canvas-Token") or request.args.get("canvas_token")
    try:
        file_bytes = _download_canvas_bytes(file_url, canvas_token)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Proxy download failed: {e}"}), 502

    # Guess content type from URL / extension
    lower = file_url.lower().split("?")[0]
    content_type = "application/octet-stream"
    if lower.endswith(".pdf"):
        content_type = "application/pdf"
    elif lower.endswith(".docx"):
        content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif lower.endswith(".pptx"):
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    elif lower.endswith(".xlsx"):
        content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    elif lower.endswith(".txt"):
        content_type = "text/plain"
    elif lower.endswith((".png", ".jpg", ".jpeg", ".webp", ".gif")):
        content_type = "image/" + lower.rsplit(".", 1)[-1].replace("jpg", "jpeg")

    from flask import Response
    return Response(
        file_bytes,
        status=200,
        mimetype=content_type,
        headers={
            "Cache-Control": "private, max-age=300",
            "Access-Control-Allow-Origin": "*",
        },
    )


@app.route("/api/ingest/canvas-url", methods=["POST"])
async def ingest_canvas_url():
    """
    Downloads a file from a given URL and ingests it.
    """
    try:
        data = request.get_json(force=True)
        assignment_id = data.get("assignment_id")
        file_url = data.get("url")
        filename = (data.get("filename") or "").lower()
        canvas_token = data.get("canvas_token")

        if not assignment_id or not file_url:
            return jsonify({"error": "assignment_id and url are required"}), 400

        rejection = _canvas_url_rejection(str(file_url))
        if rejection:
            return jsonify({"error": rejection}), 400

        headers = {}
        if canvas_token:
            headers["Authorization"] = f"Bearer {canvas_token}"

        with requests.get(file_url, headers=headers, timeout=30, stream=True) as r:
            if r.status_code != 200:
                return jsonify({"error": f"Failed to download from Canvas: {r.status_code} {r.text[:500]}"}), 400

            declared_length = r.headers.get("Content-Length")
            if declared_length and declared_length.isdigit() and int(declared_length) > MAX_CANVAS_DOWNLOAD_BYTES:
                return jsonify({"error": "File is too large to ingest (limit 25MB)"}), 400

            parts: list[bytes] = []
            total = 0
            for part in r.iter_content(chunk_size=65536):
                total += len(part)
                if total > MAX_CANVAS_DOWNLOAD_BYTES:
                    return jsonify({"error": "File is too large to ingest (limit 25MB)"}), 400
                parts.append(part)
            file_bytes = b"".join(parts)

        # Determine file type and extract text
        if filename.endswith(".pdf"):
            course_material_text = _parse_pdf(file_bytes)
        elif filename.endswith(".docx"):
            course_material_text = _parse_docx(file_bytes)
        elif filename.endswith(".pptx"):
            course_material_text = _parse_pptx(file_bytes)
        elif filename.endswith(".txt"):
            course_material_text = file_bytes.decode("utf-8", errors="ignore")
        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
            course_material_text = "[Image material submitted, skipping text extraction]"
        else:
            return jsonify({"error": f"Unsupported file type: '{filename}'."}), 400

        if not course_material_text.strip():
            return jsonify({"error": "Could not extract any text from the uploaded file"}), 400

        from grading_server.ingest_service import ingest_course_file, ingest_course_material

        course_id = data.get("course_id")
        canvas_file_id = data.get("canvas_file_id")
        updated_at = data.get("updated_at") or ""
        custom_gemini_key = request.headers.get("X-Gemini-Api-Key")

        if course_id and canvas_file_id:
            result = await ingest_course_file(
                course_id=str(course_id),
                canvas_file_id=str(canvas_file_id),
                updated_at=str(updated_at),
                filename=data.get("filename") or filename,
                course_material_text=course_material_text,
                assignment_id=assignment_id,
                api_key=custom_gemini_key,
            )
            return jsonify(result), 200

        # Legacy fallback (assignment-scoped)
        result = await ingest_course_material(
            assignment_id=assignment_id,
            course_material_text=course_material_text,
            rubric_criteria=[],
            source_name=filename,
            api_key=custom_gemini_key,
        )
        return jsonify(result), 200

    except Exception:
        traceback.print_exc()
        return jsonify({"error": "Canvas URL ingestion failed"}), 500


@app.route("/api/ingest/course-status", methods=["GET"])
async def course_ingest_status():
    course_id = request.args.get("course_id")
    if not course_id:
        return jsonify({"error": "course_id is required"}), 400
    from grading_server.ingest_service import list_course_file_status
    files = await list_course_file_status(course_id)
    return jsonify({"course_id": course_id, "files": files}), 200



@app.route("/api/ingest/<assignment_id>", methods=["DELETE"])
async def delete_ingested(assignment_id: str):
    """
    DELETE /api/ingest/<assignment_id>
    Clears all stored chunks, criteria, and exemplars for a given assignment.
    """
    try:
        if not assignment_id:
            return jsonify({"error": "assignment_id is required"}), 400

        result = await delete_assignment_data(assignment_id)
        return jsonify(result), 200

    except Exception:
        traceback.print_exc()
        return jsonify({"error": "Deletion failed"}), 500


# ── File extraction helpers ───────────────────────────────────────────────────

def _extract_pdf(file_bytes: bytes, api_key: str | None = None) -> str:
    """Extract PDF text, tables, and describe embedded charts/images via vision."""
    return _parse_pdf(file_bytes, api_key=api_key)


def _extract_docx(file_bytes: bytes, api_key: str | None = None) -> str:
    """Extract DOCX text/tables and describe embedded figures via vision."""
    return _parse_docx(file_bytes, api_key=api_key)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        lines.append("")
    return "\n".join(lines)


def _extract_image(file_bytes: bytes, mime_type: str, api_key: str | None = None) -> str:
    """Extract text/visual info from a standalone image using Gemini vision."""
    from grading_server.utils.file_parsers import _describe_visual
    desc = _describe_visual(file_bytes, mime_type, api_key=api_key)
    if desc:
        return f"[FIGURE / VISUAL - treat as present in the submission]\n{desc}"
    return ""


def _extract_ipynb(file_bytes: bytes) -> str:
    """Extract markdown, code, and text outputs from a Jupyter notebook."""
    import json

    try:
        nb = json.loads(file_bytes.decode("utf-8", errors="replace"))
    except Exception as e:
        raise ValueError(f"Invalid notebook JSON: {e}") from e

    parts: list[str] = ["=== Jupyter Notebook ==="]
    for i, cell in enumerate(nb.get("cells") or [], start=1):
        cell_type = cell.get("cell_type") or "code"
        source = cell.get("source") or ""
        if isinstance(source, list):
            source = "".join(source)
        source = str(source).strip()
        parts.append(f"\n--- Cell {i} ({cell_type}) ---")
        if source:
            parts.append(source)

        for out in cell.get("outputs") or []:
            if out.get("output_type") == "stream":
                text = out.get("text") or ""
                if isinstance(text, list):
                    text = "".join(text)
                if str(text).strip():
                    parts.append(f"[stdout]\n{text}")
            elif out.get("output_type") == "error":
                parts.append(
                    f"[error] {out.get('ename', '')}: {out.get('evalue', '')}"
                )
            else:
                data = out.get("data") or {}
                plain = data.get("text/plain") or ""
                if isinstance(plain, list):
                    plain = "".join(plain)
                if str(plain).strip():
                    parts.append(f"[output]\n{plain}")
                if data.get("image/png") or data.get("image/jpeg"):
                    parts.append("[FIGURE / VISUAL - notebook plot/image output present]")

    return "\n".join(parts).strip()


if __name__ == "__main__":
    print(f"Starting Flask server on port {FLASK_PORT}...")
    # use_reloader=False: Windows watchdog reloads kill in-flight grade/ingest
    # requests with "cannot schedule new futures after shutdown".
    app.run(host="0.0.0.0", port=FLASK_PORT, debug=True, use_reloader=False)
